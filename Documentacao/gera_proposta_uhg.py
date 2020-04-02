# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import sys
import copy
import datetime
import getopt

"""
argumentos:

gera_proposta_uhg.py --tipo_analise x --tipo_teste x (--client x || --area x) --perspectiva x --nome_resp x --prazo x

Exemplo:

gera_proposta_uhg.py --tipo_analise web --tipo_teste gray --cliente rabobank --persperctiva

--tipo_analise [web, mobile ou infra]
--tipo_teste [whitebox, graybox, blackbox]
--cliente nome_cliente
--area nome_area
--perspectiva [interno ou externo]
--nome_resp nome_responsável
--prazo prazo no formato [dd de mês de aaaa a dd de mês de aaaa]
"""

##########################################################################
#Toda vez que for necessário alterar o template							##
# deve se fazer a enumeração dos placeholders com 						##
# o seguinte código:													##
#  for shape in slide.placeholders:										##
#	print('%d %s' % (shape.placeholder_format.idx, shape.name))			##
##########################################################################

meses = ["Unknown",
         "janeiro",
         "fevereiro",
         "março",
         "abril",
         "maio",
         "junho",
         "julho",
         "agosto",
         "setembro",
         "outubro",
         "novembro",
         "dezembro"]

##################################
##	ID dos placeholders do pptx	##
##################################
ID_PLACEHOLDER_TITLE = 0
ID_PLACEHOLDER_SUBTITLE = 1
ID_PLACEHOLDER_CONTENT = 13

##############################################
## Insere os textos no slide informado por  ##
## parâmetro                                ##
##############################################
def insert_text_into_slide(slide, id, text):
    tx = slide.shapes.placeholders[id].text_frame.paragraphs[0].add_run()
    tx.text = text


##############################################################################################################
##              INÍCIO DA BRINCADEIRA                                                                       ##
##############################################################################################################
proposta = Presentation()
template = Presentation('Template_Proposta_UHG.pptx')

try:
    opts, args = getopt.getopt(sys.argv[1:], "ha:o", ["tipo-analise=", "tipo-teste=", "area=", "cliente=", "perspectiva=", "prazo=", "nome-resp="])
except getopt.GetoptError:
    print("Usage: gera_proposta_uhg.py --tipo-analise x --tipo-teste x (--cliente x || --area x) --perspectiva x --nome-resp x --prazo x")
    sys.exit(1)

for opt, arg in opts:
    if opt == '--tipo-analise':
        tipo_analise = arg
    if opt == '--tipo-teste':
        tipo_teste = arg
    if opt == '--cliente' or opt == '--area':
        cli_area = arg
    if opt == '--perspectiva':
        perspectiva = arg
    if opt == '--prazo':
        prazo = arg
    if opt == '--nome-resp':
        nome_resp = arg

##############################################
## Determina o tamanho dos slides           ##
##############################################
proposta.slide_width = 12192403
proposta.slide_height = 6851483

##############################################
## Set the slide templates                  ##
##############################################
slide_contra_capa = template.slide_masters[0].slide_layouts[0]
text_contra_capa = "São Paulo, " + str(datetime.datetime.now().day) + " de " + meses[datetime.datetime.now().month] + " de " + str(datetime.datetime.now().year) + "."
text_contra_capa2 = "À area de " + cli_area + "\nAt. " + nome_resp
text_contra_capa3 = "Prezado "+ nome_resp + ",\n\nTemos a satisfação de encaminhar à sua apreciação nossa proposta de prestação de serviço para realização de Testes de Ataque e Invasão (Pen Test).\n\nColocamo-nos a sua disposição para maiores informações que se fizerem necessárias na busca de atingirmos um melhor entendimento para cumprirmos o objetivo e escopo propostos.\n\nAtenciosamente."


slide_capa = template.slide_masters[0].slide_layouts[1]
slide_sumario = template.slide_masters[0].slide_layouts[2]
slide_obj_abrang_capa = template.slide_masters[0].slide_layouts[3]
slide_obj_abrang = template.slide_masters[0].slide_layouts[4]
text_obj = "A área de " + cli_area + " solicita a equipe de Red Team da UHG  para que seja realizado Testes de Ataque e Invasão em sua(s) Aplicação(ões) Web da perspectiva " + perspectiva + " do ambiente, na modalidade " + tipo_teste + ", com a seguinte abrangência:"
text_obj2 = "Mapear e analisar o ambiente de tecnologia corporativa voltado a internet;\nEnumerar as informações disponíveis no ambiente de Internet a partir da listagem de endereços fornecidos;\nMapear o estado das portas (TCP/UDP) de comunicação dos sistemas identificados;\nUtilizar diferentes mecanismos de exploração de falhas para ultrapassar eventuais controles de segurança e alerta;\nAnalisar vulnerabilidades na infraestrutura e/ou aplicações, definidas no escopo, e eliminação de falsos positivos;\nExplorar eventuais vulnerabilidades em outros ativos, definidos no escopo, para conseguir chegar aos alvos;\nColetar evidências das explorações e documentar, em formato de relatório, todos os pontos de falhas encontrados."

slide_metod_capa = template.slide_masters[0].slide_layouts[5]
slide_metod = template.slide_masters[0].slide_layouts[6]
slide_escop_cron_capa = template.slide_masters[0].slide_layouts[7]
slide_escop_cron = template.slide_masters[0].slide_layouts[8]
slide_gerencia_projeto_capa = template.slide_masters[0].slide_layouts[9]
slide_gerencia_projeto = template.slide_masters[0].slide_layouts[10]
slide_conf_lim_resp_capa = template.slide_masters[0].slide_layouts[11]
slide_conf_lim_resp = template.slide_masters[0].slide_layouts[12]
text_conf_lim_resp = "Confirmamos que toda e qualquer comunicação entre a equipe de Red Team da UHG e a área de " + cli_area + ", bem como quaisquer materiais ou informação desenvolvida ou recebida por nós, na forma do presente contrato, verbal ou escrita, será considerada confidencial. Assim, assumimos o compromisso de proteger as informações confidenciais contra divulgação a terceiros.\n\nA responsabilidade máxima da equipe de Red Team da UHG, de acordo com os serviços desta proposta, deve ser limitada aos serviços ou produtos do projeto, dos quais decorre a obrigação. De forma alguma deve a equipe de Red Team da UHG ser responsabilizada por perdas, danos ou despesas. Esta cláusula deverá permanecer válida além do término deste acordo.\n\nOs produtos finais elaborados pela a equipe de Red Team da UHG serão de uso exclusivo e interno da área de " + cli_area + " e não deverão ser utilizados para nenhum outro propósito.\n\nA área de " + cli_area + "deverá isentar de responsabilidade a equipe de Red Team da UHG e seu pessoal de toda e qualquer ação judicial, obrigações, custos e despesas (incluindo, sem limitações, honorários advocatícios e tempo do pessoal da equipe de Red Team da UHG envolvido), a qualquer tempo e em qualquer situação decorrente dos, ou relativa aos, serviços contidos nesta proposta. Esta cláusula deverá permanecer válida além do término deste acordo por qualquer razão.\n\nNão se constitui quebra de sigilo ou confidencialidade a simples menção a terceiros do nome e a natureza dos trabalhos prestados por meio desta proposta de serviços, desde que preservados os resultados e demais informações consideradas proprietárias."

slide_prazo_val_form_capa = template.slide_masters[0].slide_layouts[13]
slide_prazo_val_form = template.slide_masters[0].slide_layouts[14]
text_prazo_val_form="Se os termos desta proposta forem aceitos, e o conteúdo mencionado estiver de acordo com as necessidades da área de " + cli_area + ", solicitamos a devolução de uma cópia devidamente assinada, o que passará a valer como documento oficial de serviços profissionais realizados pelo prazo de " + prazo + ".\n\n\n\nDe acordo:" + cli_area + " - " + nome_resp
slide_ultimo = template.slide_masters[0].slide_layouts[15]

###############################
## Slide capa                ##
###############################
slide_capa_rel = proposta.slides.add_slide(slide_capa)
insert_text_into_slide(slide_capa_rel, ID_PLACEHOLDER_CONTENT - 3, tipo_analise + ' - ' + tipo_teste)

###############################
## Slide contra capa         ##
###############################
slide_contra_capa_proposta = proposta.slides.add_slide(slide_contra_capa)
insert_text_into_slide(slide_contra_capa_proposta, ID_PLACEHOLDER_CONTENT-3, text_contra_capa)
insert_text_into_slide(slide_contra_capa_proposta, ID_PLACEHOLDER_CONTENT-2, text_contra_capa2)
insert_text_into_slide(slide_contra_capa_proposta, ID_PLACEHOLDER_CONTENT-1, text_contra_capa3)

###############################
## Slide sumário             ##
###############################
proposta.slides.add_slide(slide_sumario)

###############################
## Slide objetivo e          ##
## metodologia               ##
###############################
proposta.slides.add_slide(slide_obj_abrang_capa)

###############################
## Slide metodologia e       ##
## escopo                    ##
###############################
slide_obj_abrang_proposta = proposta.slides.add_slide(slide_obj_abrang)
insert_text_into_slide(slide_obj_abrang_proposta, ID_PLACEHOLDER_CONTENT + 1, text_obj)
insert_text_into_slide(slide_obj_abrang_proposta, ID_PLACEHOLDER_CONTENT, text_obj2)

proposta.slides.add_slide(slide_metod_capa)

##################################
## Slide plano de ação 			##
## alimentando os títulos de    ##
## cada ponto no slide		    ##
##################################
proposta.slides.add_slide(slide_metod)
proposta.slides.add_slide(slide_escop_cron_capa)

proposta.slides.add_slide(slide_gerencia_projeto_capa)

proposta.slides.add_slide(slide_gerencia_projeto)
##########################################################################
#Adiciona slide da descrição do resultado obtido do ponto com			##
#título, descrição, código da vulnerabilidade, criticidade,				##
#impact, recomendação e referências										##
##########################################################################
proposta.slides.add_slide(slide_conf_lim_resp_capa)
slide_conf_lim_resp_proposta = proposta.slides.add_slide(slide_conf_lim_resp)
insert_text_into_slide(slide_conf_lim_resp_proposta, ID_PLACEHOLDER_CONTENT, text_conf_lim_resp)


proposta.slides.add_slide(slide_prazo_val_form_capa)
slide_prazo_val_form_proposta = proposta.slides.add_slide(slide_prazo_val_form)
insert_text_into_slide(slide_prazo_val_form_proposta, ID_PLACEHOLDER_CONTENT, text_prazo_val_form)
proposta.slides.add_slide(slide_ultimo)

proposta.save('proposta.pptx')
